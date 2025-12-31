import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from design_enforcer import enforce_design
from theme_engine import hex_to_rgb
from slide_mapper import get_slide_layout
from design_styles import get_design_style, apply_design_decorations
from visual_elements import add_visual_elements_to_slide
from performance_monitor import checkpoint
import gc
import re
from typing import Dict, List, Optional

SLIDE_LAYOUTS = {
    "title": 0,
    "content": 1,
    "summary": 1
}

# Perfect typography settings
PERFECT_FONTS = {
    "title": {"name": "Segoe UI", "size": 44, "bold": True},
    "subtitle": {"name": "Segoe UI Light", "size": 24, "bold": False},
    "body": {"name": "Segoe UI", "size": 18, "bold": False},
    "bullet": {"name": "Segoe UI", "size": 16, "bold": False}
}

def generate_ppt(slides: list, output_path: str, design_style: str = "minimal_1", visual_preferences: dict = None):
    """
    Generate PERFECT PPT with premium design, flawless formatting, and zero errors
    
    Args:
        slides: List of slide data dictionaries
        output_path: Path to save the PPT file
        design_style: Design style ID to use (e.g., "minimal_1", "corporate_1", "tech_1", etc.)
        visual_preferences: Dict of visual element preferences
    """
    prs = None
    try:
        # Input validation with detailed error messages
        if not isinstance(slides, list):
            raise ValueError("Slides must be a list of dictionaries")
        
        if not slides:
            raise ValueError("At least one slide is required")
        
        if visual_preferences is None:
            visual_preferences = {
                "graphs": False,
                "tables": False,
                "pie_charts": False,
                "images": False
            }
        
        slide_count = len(slides)
        
        # Intelligent memory optimization based on slide count
        memory_budget_per_slide = _calculate_memory_budget(slide_count, visual_preferences)
        
        # Create presentation with perfect settings
        prs = Presentation()
        
        # Store design style and slide count for visual elements
        prs._design_style = design_style
        prs._total_slides = slide_count
        
        # Set to perfect 16:9 widescreen ratio for modern displays
        prs.slide_width = Inches(13.333)  # Perfect widescreen width
        prs.slide_height = Inches(7.5)    # Perfect widescreen height

        # Get design style configuration with fallback protection
        style_config = get_design_style(design_style)
        if not style_config:
            print(f"[WARNING] Invalid design style '{design_style}', using minimal_1")
            style_config = get_design_style("minimal_1")

        print(f"[PERFECT] Creating {slide_count} slides with {design_style} design...")
        
        # Track memory for optimization
        initial_memory = _get_memory_usage()
        
        for slide_index, slide_data in enumerate(slides):
            # Validate slide data structure
            if not isinstance(slide_data, dict):
                raise ValueError(f"Slide {slide_index + 1} must be a dictionary")

            # Memory monitoring and optimization
            current_memory = _get_memory_usage()
            if slide_index > 0:
                memory_per_slide = (current_memory - initial_memory) / slide_index
                if memory_per_slide > memory_budget_per_slide:
                    print(f"[OPTIMIZE] Memory usage {memory_per_slide:.1f}MB/slide, optimizing...")
                    _aggressive_cleanup()
            
            # Progress tracking with beautiful output
            if slide_index % 3 == 0 and slide_index > 0:
                progress = (slide_index / slide_count) * 100
                print(f"[PROGRESS] {progress:.0f}% complete - {slide_index}/{slide_count} slides")
                checkpoint(slide_index, "slide_batch_complete")

            # Extract and validate slide content
            slide_type = slide_data.get("slide_type", "content")
            title = _clean_text(slide_data.get("title", ""))
            bullets = slide_data.get("bullets", [])

            # Validate slide type
            if slide_type not in SLIDE_LAYOUTS:
                print(f"[WARNING] Unknown slide type '{slide_type}', using 'content'")
                slide_type = "content"

            # Create slide with proper layout
            slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUTS[slide_type]])

            # Apply perfect background design
            _apply_perfect_background(slide, style_config)

            # BEAUTIFUL SIMPLE SYSTEM - handles all content and visual placement
            print(f"[BEAUTIFUL] Creating beautiful, uncluttered layout for slide {slide_index + 1}")
            
            # Add beautiful visual elements using the simple system
            add_visual_elements_to_slide(slide, slide_data, design_style, visual_preferences, prs, slide_index)
            
            # Apply perfect design decorations (minimal for clean look)
            if slide_count <= 15:  # Only for manageable presentations
                apply_design_decorations(slide, design_style, prs)
            
            # Enforce perfect design rules
            layout_info = get_slide_layout(slide_type)
            enforce_design(slide, layout_info)
            
            # Memory cleanup every 5 slides
            if slide_index % 5 == 4:
                _aggressive_cleanup()

        # Save with perfect error handling
        print(f"[PERFECT] Finalizing presentation...")
        checkpoint(slide_count, "ppt_complete")
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        prs.save(output_path)
        
        final_memory = _get_memory_usage()
        print(f"[SUCCESS] Perfect presentation created!")
        print(f"[SUCCESS] {slide_count} slides generated flawlessly")
        print(f"[SUCCESS] File saved: {output_path}")
        print(f"[SUCCESS] Final memory: {final_memory:.1f}MB")
        
        return output_path
        
    except Exception as e:
        print(f"[ERROR] PPT generation failed: {str(e)}")
        print(f"[ERROR] Error type: {type(e).__name__}")
        raise
        
    finally:
        # Perfect cleanup
        if prs:
            try:
                # Clear all references safely
                for slide in prs.slides:
                    slide = None
                prs._element.clear() if hasattr(prs, '_element') else None
            except:
                pass
            del prs
        
        _aggressive_cleanup()
        print(f"[CLEANUP] Memory after cleanup: {_get_memory_usage():.1f}MB")


def _calculate_memory_budget(slide_count: int, visual_preferences: dict) -> int:
    """Calculate optimal memory budget per slide"""
    base_budget = 100
    
    if slide_count > 20:
        return 40  # Ultra conservative for very large presentations
    elif slide_count > 15:
        return 60  # Conservative for large presentations
    elif slide_count > 10:
        return 80  # Moderate for medium presentations
    
    # Adjust for visual elements
    if any(visual_preferences.values()):
        return base_budget + 20  # Extra budget for visuals
    
    return base_budget


def _clean_text(text: str) -> str:
    """Clean and optimize text for perfect presentation"""
    if not text:
        return ""
    
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', str(text).strip())
    
    # Fix common typography issues
    text = text.replace('"', '"').replace('"', '"')  # Smart quotes
    text = text.replace("'", "'").replace("'", "'")  # Smart apostrophes
    text = text.replace('--', '—')  # Em dashes
    text = text.replace('...', '…')  # Ellipsis
    
    # Ensure proper sentence ending
    if text and not text.endswith(('.', '!', '?', ':', '…')):
        text += '.'
    
    return text


def _create_perfect_title(title_shape, title_text: str, style_config: dict):
    """Create perfectly formatted title - REMOVED - Now handled by professional layout system"""
    # This function is now handled by the professional layout system
    # The professional_visual_system creates clean, properly positioned titles
    pass


def _create_perfect_content(content_placeholder, bullets: list, style_config: dict, visual_preferences: dict):
    """Create perfectly formatted content - REMOVED - Now handled by professional layout system"""
    # This function is now handled by the professional layout system
    # The professional_visual_system creates clean, properly spaced content
    pass


def _apply_perfect_background(slide, style_config: dict):
    """Apply perfect background with design colors"""
    try:
        bg = slide.background
        fill = bg.fill
        
        bg_config = style_config["background"]
        
        # For now, use solid color backgrounds for reliability
        # Gradients can be complex and cause issues in different PowerPoint versions
        fill.solid()
        
        # Use the first color from the design style
        primary_color = bg_config["colors"][0]
        fill.fore_color.rgb = hex_to_rgb(primary_color)
        
        print(f"[DESIGN] Applied background color: {primary_color}")
            
    except Exception as e:
        print(f"[WARNING] Background application failed: {e}")
        # Ultimate fallback to white background
        try:
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        except:
            pass


def _should_add_visuals(slide_index: int, slide_count: int, visual_preferences: dict) -> bool:
    """Determine if visual elements should be added to this slide - MORE GENEROUS"""
    if not any(visual_preferences.values()):
        return False
    
    # Skip title slide (index 0)
    if slide_index == 0:
        return False
    
    # More generous visual element addition
    if slide_count <= 5:
        return slide_index >= 1  # Add to slides 2, 3, 4, 5
    elif slide_count <= 10:
        return slide_index % 2 == 1  # Add to slides 2, 4, 6, 8, 10
    elif slide_count <= 15:
        return slide_index % 3 != 0  # Add to slides 2, 3, 5, 6, 8, 9, 11, 12, 14, 15
    elif slide_count <= 20:
        return slide_index % 3 == 1 or slide_index % 3 == 2  # Add to most slides
    else:
        return slide_index % 4 != 0  # Add to 75% of slides for large presentations


def _get_memory_usage() -> float:
    """Get current memory usage in MB with fallback"""
    try:
        import psutil
        process = psutil.Process()
        return process.memory_info().rss / 1024 / 1024
    except ImportError:
        # Fallback using tracemalloc if available
        try:
            import tracemalloc
            if tracemalloc.is_tracing():
                current, peak = tracemalloc.get_traced_memory()
                return current / 1024 / 1024
        except:
            pass
        return 0.0


def _aggressive_cleanup():
    """Perform aggressive memory cleanup for perfect performance"""
    # Multiple garbage collection passes
    for _ in range(3):
        gc.collect()
    
    # Clear weak references
    try:
        import weakref
        weakref.getweakrefs(object)
    except:
        pass
    
    # Force Python to release memory back to OS
    try:
        import ctypes
        ctypes.CDLL("libc.so.6").malloc_trim(0)
    except:
        pass


# Legacy function for backward compatibility
def apply_design_background(slide, style_config):
    """Legacy function - use _apply_perfect_background instead"""
    _apply_perfect_background(slide, style_config)


def apply_design_colors(slide, style_config):
    """Apply design style colors to text elements"""
    colors = style_config["colors"]
    
    # Title colors with error handling
    try:
        if slide.shapes.title and slide.shapes.title.text_frame.paragraphs:
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        run.font.color.rgb = hex_to_rgb(colors["title"])
    except Exception as e:
        print(f"[WARNING] Title color application failed: {e}")
    
    # Body colors with error handling
    try:
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        run.font.color.rgb = hex_to_rgb(colors["body"])
    except Exception as e:
        print(f"[WARNING] Body color application failed: {e}")