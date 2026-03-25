# design_engine.py
# Complete Design Engine - Each design has unique layout, decorations, typography

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ═══════════════════════════════════════════════════════════════════
#  DESIGN CATALOGUE
#  Each entry defines: background, title_style, body_style,
#  accent_bar, corner_shape, divider, font_title, font_body
# ═══════════════════════════════════════════════════════════════════

DESIGNS = {

    # ── MINIMAL ─────────────────────────────────────────────────────
    'minimal_1': {
        'name': 'Clean White',
        'category': 'Minimal',
        'bg': '#FFFFFF',
        'title_color': '#1A1A2E',
        'body_color':  '#3D3D3D',
        'accent':      '#4361EE',
        'font_title':  'Helvetica Neue',
        'font_body':   'Helvetica Neue',
        # Thin left accent bar
        'left_bar': {'color': '#4361EE', 'width': Inches(0.07)},
        'top_line': None,
        'corner': None,
        'bottom_strip': None,
    },
    'minimal_2': {
        'name': 'Warm Paper',
        'category': 'Minimal',
        'bg': '#FAF7F2',
        'title_color': '#2D2D2D',
        'body_color':  '#555555',
        'accent':      '#C9A84C',
        'font_title':  'Georgia',
        'font_body':   'Calibri',
        # Warm gold top border
        'left_bar': None,
        'top_line': {'color': '#C9A84C', 'height': Inches(0.06)},
        'corner': None,
        'bottom_strip': None,
    },
    'minimal_3': {
        'name': 'Ink Black',
        'category': 'Minimal',
        'bg': '#F5F5F5',
        'title_color': '#111111',
        'body_color':  '#444444',
        'accent':      '#111111',
        'font_title':  'Calibri Light',
        'font_body':   'Calibri',
        # Bold bottom strip
        'left_bar': None,
        'top_line': None,
        'corner': None,
        'bottom_strip': {'color': '#111111', 'height': Inches(0.12)},
    },

    # ── CORPORATE ───────────────────────────────────────────────────
    'corporate_1': {
        'name': 'Executive Navy',
        'category': 'Corporate',
        'bg': '#0D1B2A',
        'title_color': '#FFFFFF',
        'body_color':  '#CBD5E1',
        'accent':      '#F59E0B',
        'font_title':  'Calibri',
        'font_body':   'Calibri',
        # Gold left bar + gold top line
        'left_bar': {'color': '#F59E0B', 'width': Inches(0.1)},
        'top_line': {'color': '#F59E0B', 'height': Inches(0.04)},
        'corner': None,
        'bottom_strip': None,
    },
    'corporate_2': {
        'name': 'Steel Gray',
        'category': 'Corporate',
        'bg': '#1C1C1E',
        'title_color': '#FFFFFF',
        'body_color':  '#AEAEB2',
        'accent':      '#0A84FF',
        'font_title':  'Segoe UI',
        'font_body':   'Segoe UI',
        # Blue right corner block
        'left_bar': None,
        'top_line': None,
        'corner': {'color': '#0A84FF', 'size': Inches(1.2), 'position': 'top-right'},
        'bottom_strip': None,
    },
    'corporate_3': {
        'name': 'Forest Executive',
        'category': 'Corporate',
        'bg': '#1B2B1F',
        'title_color': '#FFFFFF',
        'body_color':  '#D1FAE5',
        'accent':      '#34D399',
        'font_title':  'Calibri',
        'font_body':   'Calibri',
        # Green bottom strip
        'left_bar': None,
        'top_line': None,
        'corner': None,
        'bottom_strip': {'color': '#34D399', 'height': Inches(0.1)},
    },

    # ── TECH ────────────────────────────────────────────────────────
    'tech_1': {
        'name': 'Cyber Dark',
        'category': 'Tech',
        'bg': '#0A0E1A',
        'title_color': '#00F5FF',
        'body_color':  '#B0C4DE',
        'accent':      '#00F5FF',
        'font_title':  'Consolas',
        'font_body':   'Calibri',
        # Cyan left bar
        'left_bar': {'color': '#00F5FF', 'width': Inches(0.06)},
        'top_line': {'color': '#00F5FF', 'height': Inches(0.03)},
        'corner': None,
        'bottom_strip': None,
    },
    'tech_2': {
        'name': 'Matrix Green',
        'category': 'Tech',
        'bg': '#0D1117',
        'title_color': '#39FF14',
        'body_color':  '#C9D1D9',
        'accent':      '#39FF14',
        'font_title':  'Consolas',
        'font_body':   'Calibri',
        # Green top line
        'left_bar': None,
        'top_line': {'color': '#39FF14', 'height': Inches(0.04)},
        'corner': None,
        'bottom_strip': {'color': '#39FF14', 'height': Inches(0.04)},
    },
    'tech_3': {
        'name': 'Purple Neon',
        'category': 'Tech',
        'bg': '#13001E',
        'title_color': '#E040FB',
        'body_color':  '#CE93D8',
        'accent':      '#E040FB',
        'font_title':  'Segoe UI',
        'font_body':   'Calibri',
        # Purple corner
        'left_bar': None,
        'top_line': None,
        'corner': {'color': '#E040FB', 'size': Inches(1.5), 'position': 'top-left'},
        'bottom_strip': None,
    },

    # ── CREATIVE ────────────────────────────────────────────────────
    'creative_1': {
        'name': 'Coral Burst',
        'category': 'Creative',
        'bg': '#FF6B6B',
        'title_color': '#FFFFFF',
        'body_color':  '#FFE8E8',
        'accent':      '#FFE66D',
        'font_title':  'Calibri Light',
        'font_body':   'Calibri',
        # Yellow bottom strip
        'left_bar': None,
        'top_line': None,
        'corner': None,
        'bottom_strip': {'color': '#FFE66D', 'height': Inches(0.15)},
    },
    'creative_2': {
        'name': 'Ocean Depth',
        'category': 'Creative',
        'bg': '#1A1A2E',
        'title_color': '#E94560',
        'body_color':  '#A8DADC',
        'accent':      '#E94560',
        'font_title':  'Calibri Light',
        'font_body':   'Calibri',
        # Red left bar
        'left_bar': {'color': '#E94560', 'width': Inches(0.12)},
        'top_line': None,
        'corner': None,
        'bottom_strip': None,
    },
    'creative_3': {
        'name': 'Mango Splash',
        'category': 'Creative',
        'bg': '#FFFBF0',
        'title_color': '#FF6D00',
        'body_color':  '#5D4037',
        'accent':      '#FF6D00',
        'font_title':  'Georgia',
        'font_body':   'Calibri',
        # Orange top line
        'left_bar': None,
        'top_line': {'color': '#FF6D00', 'height': Inches(0.08)},
        'corner': None,
        'bottom_strip': None,
    },

    # ── ACADEMIC ────────────────────────────────────────────────────
    'academic_1': {
        'name': 'Scholar Green',
        'category': 'Academic',
        'bg': '#F0F7F4',
        'title_color': '#1B4332',
        'body_color':  '#2D6A4F',
        'accent':      '#1B4332',
        'font_title':  'Georgia',
        'font_body':   'Calibri',
        # Dark green left bar
        'left_bar': {'color': '#1B4332', 'width': Inches(0.08)},
        'top_line': None,
        'corner': None,
        'bottom_strip': None,
    },
    'academic_2': {
        'name': 'Oxford Classic',
        'category': 'Academic',
        'bg': '#F8F4EF',
        'title_color': '#002147',
        'body_color':  '#3D3D3D',
        'accent':      '#002147',
        'font_title':  'Georgia',
        'font_body':   'Calibri',
        # Navy top + bottom lines (classic book style)
        'left_bar': None,
        'top_line': {'color': '#002147', 'height': Inches(0.05)},
        'corner': None,
        'bottom_strip': {'color': '#002147', 'height': Inches(0.05)},
    },
    'academic_3': {
        'name': 'Burgundy Scholar',
        'category': 'Academic',
        'bg': '#FDF6F0',
        'title_color': '#6B0F1A',
        'body_color':  '#4A1020',
        'accent':      '#6B0F1A',
        'font_title':  'Georgia',
        'font_body':   'Calibri',
        # Burgundy left bar
        'left_bar': {'color': '#6B0F1A', 'width': Inches(0.09)},
        'top_line': None,
        'corner': None,
        'bottom_strip': None,
    },

    # ── MODERN ──────────────────────────────────────────────────────
    'modern_1': {
        'name': 'Violet Gradient',
        'category': 'Modern',
        'bg': '#7C3AED',
        'title_color': '#FFFFFF',
        'body_color':  '#EDE9FE',
        'accent':      '#FCD34D',
        'font_title':  'Segoe UI Light',
        'font_body':   'Segoe UI',
        # Yellow bottom strip
        'left_bar': None,
        'top_line': None,
        'corner': None,
        'bottom_strip': {'color': '#FCD34D', 'height': Inches(0.1)},
    },
    'modern_2': {
        'name': 'Rose Gold',
        'category': 'Modern',
        'bg': '#FFF0F3',
        'title_color': '#C9184A',
        'body_color':  '#590D22',
        'accent':      '#C9184A',
        'font_title':  'Calibri Light',
        'font_body':   'Calibri',
        # Rose left bar
        'left_bar': {'color': '#C9184A', 'width': Inches(0.07)},
        'top_line': None,
        'corner': None,
        'bottom_strip': None,
    },
    'modern_3': {
        'name': 'Midnight Teal',
        'category': 'Modern',
        'bg': '#0D3B38',
        'title_color': '#FFFFFF',
        'body_color':  '#B2DFDB',
        'accent':      '#4DB6AC',
        'font_title':  'Segoe UI',
        'font_body':   'Calibri',
        # Teal corner
        'left_bar': None,
        'top_line': {'color': '#4DB6AC', 'height': Inches(0.05)},
        'corner': {'color': '#4DB6AC', 'size': Inches(1.0), 'position': 'bottom-right'},
        'bottom_strip': None,
    },
}


# ═══════════════════════════════════════════════════════════════════
#  APPLY DESIGN TO SLIDE
# ═══════════════════════════════════════════════════════════════════

def apply_design(slide, prs, design_id: str):
    """Apply full design to a slide - background + decorative elements."""
    d = DESIGNS.get(design_id, DESIGNS['minimal_1'])

    sw = prs.slide_width
    sh = prs.slide_height

    # 1. Background
    _set_background(slide, d['bg'])

    # 2. Left accent bar
    if d.get('left_bar'):
        bar = d['left_bar']
        _add_rect(slide, 0, 0, bar['width'], sh, bar['color'])

    # 3. Top line
    if d.get('top_line'):
        tl = d['top_line']
        _add_rect(slide, 0, 0, sw, tl['height'], tl['color'])

    # 4. Bottom strip
    if d.get('bottom_strip'):
        bs = d['bottom_strip']
        _add_rect(slide, 0, sh - bs['height'], sw, bs['height'], bs['color'])

    # 5. Corner triangle / block
    if d.get('corner'):
        c = d['corner']
        _add_corner(slide, sw, sh, c['color'], c['size'], c['position'])


def get_design(design_id: str) -> dict:
    """Return design config dict (compatible with existing code)."""
    d = DESIGNS.get(design_id, DESIGNS['minimal_1'])
    return {
        'name': d['name'],
        'colors': {
            'title':  d['title_color'],
            'body':   d['body_color'],
            'accent': d['accent'],
        },
        'background': {
            'type':   'solid',
            'colors': [d['bg']],
        },
        'font_title': d.get('font_title', 'Calibri'),
        'font_body':  d.get('font_body',  'Calibri'),
    }


def list_designs():
    """Return list of all designs grouped by category."""
    result = {}
    for did, d in DESIGNS.items():
        cat = d['category']
        result.setdefault(cat, []).append({'id': did, 'name': d['name']})
    return result


# ═══════════════════════════════════════════════════════════════════
#  INTERNAL HELPERS
# ═══════════════════════════════════════════════════════════════════

def _set_background(slide, hex_color: str):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(hex_color)


def _add_rect(slide, left, top, width, height, hex_color: str):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(hex_color)
    shape.line.fill.background()   # no border
    return shape


def _add_corner(slide, sw, sh, hex_color: str, size, position: str):
    """Add a right-triangle corner accent."""
    if position == 'top-right':
        left = sw - size
        top  = 0
    elif position == 'top-left':
        left = 0
        top  = 0
    elif position == 'bottom-right':
        left = sw - size
        top  = sh - size
    else:  # bottom-left
        left = 0
        top  = sh - size

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(hex_color)
    shape.line.fill.background()

    # Rotate to correct corner orientation
    rotations = {
        'top-right':    270,
        'top-left':     180,
        'bottom-right':   0,
        'bottom-left':   90,
    }
    shape.rotation = rotations.get(position, 0)
    return shape
