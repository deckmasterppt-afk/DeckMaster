# design_enforcer.py
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Professional font hierarchy - modern, clean, widely available
FONT_CONFIG = {
    "title": {
        "name": "Segoe UI",  # Modern, professional, widely available
        "size": Pt(44),
        "bold": True,
        "color": RGBColor(25, 25, 25)  # Near-black for better contrast
    },
    "subtitle": {
        "name": "Segoe UI Semibold",
        "size": Pt(32),
        "bold": True,
        "color": RGBColor(50, 50, 50)
    },
    "body": {
        "name": "Segoe UI",  # Clean, readable body font
        "size": Pt(18),
        "bold": False,
        "color": RGBColor(30, 30, 30)
    },
    "bullet": {
        "name": "Segoe UI",
        "size": Pt(20),
        "bold": False,
        "color": RGBColor(40, 40, 40),
        "line_spacing": 1.2
    }
}

# Professional spacing standards
SPACING_CONFIG = {
    "title_top": Inches(0.5),
    "title_bottom": Inches(0.3),
    "content_margin": Inches(0.8),
    "bullet_spacing": Pt(12),
    "paragraph_spacing": Pt(8)
}

def enforce_design(slide, layout_info):
    """
    Apply consistent fonts, spacing, hierarchy, and contrast
    Enhanced with professional typography and spacing
    """
    # Title formatting with improved hierarchy - FIX OVERFLOW
    if "title" in layout_info["fields"] and slide.shapes.title:
        title = slide.shapes.title
        title_frame = title.text_frame
        
        # Set title position and alignment - constrain to top area
        title.top = SPACING_CONFIG["title_top"]
        title.height = Inches(1.0)  # Limit title height
        title.left = Inches(0.5)
        title.width = Inches(9)  # Full width minus margins
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.margin_left = Inches(0.2)
        title_frame.margin_right = Inches(0.2)
        title_frame.word_wrap = True
        title_frame.auto_size = None  # Control size manually
        
        # Apply professional title styling with adaptive size
        for p in title_frame.paragraphs:
            # Truncate very long titles
            if len(p.text) > 80:
                p.text = p.text[:77] + "..."
            
            p.font.name = FONT_CONFIG["title"]["name"]
            # Adaptive title size based on length
            if len(p.text) > 50:
                p.font.size = Pt(36)  # Smaller for long titles
            else:
                p.font.size = FONT_CONFIG["title"]["size"]
            p.font.bold = FONT_CONFIG["title"]["bold"]
            p.font.color.rgb = FONT_CONFIG["title"]["color"]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)  # Reduced spacing
            p.line_spacing = 1.0

    # Body/bullet formatting with improved readability - FIX OVERFLOW
    if "bullets" in layout_info["fields"] and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # CRITICAL: Set proper boundaries to prevent overflow
        # Calculate available height (slide height - title area - margins)
        slide_height = Inches(5.625)  # Standard slide height
        title_bottom = SPACING_CONFIG["title_top"] + Inches(1.0)  # Title takes ~1 inch
        available_height = slide_height - title_bottom - Inches(0.5)  # Bottom margin
        
        # Set content position and size to fit on slide
        body.top = title_bottom + Inches(0.2)  # Start below title
        body.height = available_height  # Constrain height
        body.left = SPACING_CONFIG["content_margin"]
        body.width = Inches(10) - (SPACING_CONFIG["content_margin"] * 2)  # Full width minus margins
        
        # Set content margins for better whitespace
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.word_wrap = True
        tf.auto_size = None  # Disable auto-size to control boundaries
        
        # Limit number of bullets to fit on slide (max 6-7 bullets)
        max_bullets = 6
        paragraphs_to_style = tf.paragraphs[:max_bullets] if len(tf.paragraphs) > max_bullets else tf.paragraphs
        
        # Apply professional body styling with adaptive sizing
        for idx, p in enumerate(paragraphs_to_style):
            # Level-based styling for hierarchy
            level = p.level if hasattr(p, 'level') else 0
            
            # Adaptive font size based on number of bullets
            num_bullets = len(paragraphs_to_style)
            if num_bullets > 5:
                base_size = Pt(16)  # Smaller if many bullets
            elif num_bullets > 3:
                base_size = Pt(18)
            else:
                base_size = FONT_CONFIG["bullet"]["size"]
            
            if level == 0:
                # Main bullet points
                p.font.name = FONT_CONFIG["bullet"]["name"]
                p.font.size = base_size
                p.font.bold = False
                p.font.color.rgb = FONT_CONFIG["bullet"]["color"]
            else:
                # Sub-bullets (slightly smaller)
                p.font.name = FONT_CONFIG["body"]["name"]
                p.font.size = Pt(max(14, base_size.pt - 2))  # 2pt smaller
                p.font.color.rgb = RGBColor(60, 60, 60)
            
            # Reduce spacing if many bullets
            spacing = Pt(8) if num_bullets > 5 else SPACING_CONFIG["bullet_spacing"]
            p.space_after = spacing
            p.line_spacing = 1.15  # Tighter line spacing
            p.alignment = PP_ALIGN.LEFT

    # Ensure color contrast for all text shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                # Only set default if no color is explicitly set
                if not hasattr(p.font.color, 'rgb') or p.font.color.rgb is None:
                    p.font.color.rgb = RGBColor(30, 30, 30)  # Dark gray for readability
