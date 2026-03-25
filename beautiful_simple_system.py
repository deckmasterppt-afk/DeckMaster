# beautiful_simple_system.py
# Design-aware presentation layout system
# Every margin, font size, and visual zone adapts to the active design

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io


def _hex(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class DesignContext:
    """
    Computes safe content zones based on the active design's decorations.
    Every slide element is positioned INSIDE these safe zones.
    """

    def __init__(self, design_id: str):
        from design_engine import DESIGNS
        d = DESIGNS.get(design_id, DESIGNS['minimal_1'])

        self.sw = Inches(13.333)
        self.sh = Inches(7.5)

        # Base padding
        pad = Inches(0.55)

        # Left offset: push right if there's a left bar
        lb = d.get('left_bar')
        self.x0 = (lb['width'] + Inches(0.25)) if lb else pad

        # Top offset: push down if there's a top line
        tl = d.get('top_line')
        self.y0 = (tl['height'] + Inches(0.25)) if tl else pad

        # Bottom offset: shrink if there's a bottom strip
        bs = d.get('bottom_strip')
        self.y1 = self.sh - ((bs['height'] + Inches(0.2)) if bs else pad)

        # Right offset: shrink if there's a top-right corner
        c = d.get('corner')
        if c and c.get('position') == 'top-right':
            self.x1 = self.sw - c['size'] - Inches(0.2)
        else:
            self.x1 = self.sw - pad

        # Usable dimensions
        self.w = self.x1 - self.x0
        self.h = self.y1 - self.y0

        # Design personality flags
        self.category   = d.get('category', 'Minimal')
        self.font_title = d.get('font_title', 'Calibri')
        self.font_body  = d.get('font_body',  'Calibri')
        self.title_color = d.get('title_color', '#1A1A2E')
        self.body_color  = d.get('body_color',  '#3D3D3D')
        self.accent      = d.get('accent',       '#4361EE')
        self.bg          = d.get('bg',           '#FFFFFF')

        # Is background dark? → affects chart/table colour theme
        r, g, b = int(self.bg[1:3], 16), int(self.bg[3:5], 16), int(self.bg[5:7], 16)
        self.dark_bg = (r * 0.299 + g * 0.587 + b * 0.114) < 128

        # Title alignment: centred for creative/modern, left for others
        self.title_align = (PP_ALIGN.CENTER
                            if self.category in ('Creative', 'Modern')
                            else PP_ALIGN.LEFT)

        # Title size: slightly larger for minimal (more whitespace), smaller for tech
        self.title_pt = {'Minimal': 32, 'Academic': 30, 'Corporate': 30,
                         'Tech': 28, 'Creative': 34, 'Modern': 32
                         }.get(self.category, 30)

        # Body size base
        self.body_pt_base = {'Minimal': 19, 'Academic': 18, 'Corporate': 18,
                             'Tech': 17, 'Creative': 18, 'Modern': 18
                             }.get(self.category, 18)

        # Line spacing
        self.line_spacing = 1.4 if self.category in ('Academic', 'Minimal') else 1.3

        # Divider line under title?
        self.show_divider = self.category in ('Academic', 'Corporate', 'Minimal')


class BeautifulSimpleSystem:
    def __init__(self):
        self.slide_width  = Inches(13.333)
        self.slide_height = Inches(7.5)

    # ─────────────────────────────────────────────
    # PUBLIC ENTRY POINT
    # ─────────────────────────────────────────────

    def create_beautiful_slide(self, slide, slide_data, slide_index,
                               total_slides, design_style='minimal_1'):
        self.current_design_style = design_style
        self._clear_slide(slide)

        is_last = (slide_index == total_slides - 1)
        ctx = DesignContext(design_style)

        if slide_index == 0:
            return self._title_slide(slide, slide_data, ctx)
        elif is_last:
            return self._thankyou_slide(slide, slide_data, ctx)
        else:
            return self._content_slide(slide, slide_data, ctx)

    # ─────────────────────────────────────────────
    # TITLE SLIDE
    # ─────────────────────────────────────────────

    def _title_slide(self, slide, slide_data, ctx: DesignContext):
        cx = self.slide_width  / 2
        cy = self.slide_height / 2

        # Main title
        w, h = Inches(10), Inches(2.2)
        shape = slide.shapes.add_textbox(cx - w/2, cy - h - Inches(0.2), w, h)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = slide_data.get('title', 'Presentation Title')
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = ctx.font_title
        run.font.size = Pt(54)
        run.font.bold = ctx.category in ('Corporate',)
        run.font.color.rgb = _hex(ctx.title_color)

        # Thin accent divider
        dw, dh = Inches(3), Inches(0.04)
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx - dw/2, cy + Inches(0.05), dw, dh)
        div.fill.solid()
        div.fill.fore_color.rgb = _hex(ctx.accent)
        div.line.fill.background()

        # Subtitle
        bullets = slide_data.get('bullets', [])
        if bullets:
            sw, sh = Inches(9), Inches(1.2)
            sub = slide.shapes.add_textbox(cx - sw/2, cy + Inches(0.3), sw, sh)
            stf = sub.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = str(bullets[0]).lstrip('•').strip()
            sp.alignment = PP_ALIGN.CENTER
            sr = sp.runs[0]
            sr.font.name = ctx.font_body
            sr.font.size = Pt(22)
            sr.font.italic = ctx.category in ('Academic', 'Minimal')
            sr.font.color.rgb = _hex(ctx.body_color)

        return 'title'

    # ─────────────────────────────────────────────
    # THANK YOU SLIDE
    # ─────────────────────────────────────────────

    def _thankyou_slide(self, slide, slide_data, ctx: DesignContext):
        cx = self.slide_width  / 2
        cy = self.slide_height / 2

        # "Thank You" – elegant script font
        w, h = Inches(10), Inches(2.5)
        shape = slide.shapes.add_textbox(cx - w/2, cy - h/2 - Inches(0.4), w, h)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = slide_data.get('title', 'Thank You')
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        # Script font for elegance; fall back gracefully in PowerPoint
        run.font.name = 'Segoe Script'
        run.font.size = Pt(68)
        run.font.bold = False
        run.font.color.rgb = _hex(ctx.title_color)

        # Accent divider
        dw, dh = Inches(2.5), Inches(0.04)
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx - dw/2, cy + Inches(0.25), dw, dh)
        div.fill.solid()
        div.fill.fore_color.rgb = _hex(ctx.accent)
        div.line.fill.background()

        # Subtitle
        bullets = slide_data.get('bullets', [])
        subtitle = (str(bullets[0]).lstrip('•').strip()
                    if bullets else 'We appreciate your time and attention')
        sw, sh = Inches(9), Inches(1)
        sub = slide.shapes.add_textbox(cx - sw/2, cy + Inches(0.5), sw, sh)
        stf = sub.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.runs[0]
        sr.font.name = ctx.font_body
        sr.font.size = Pt(20)
        sr.font.italic = True
        sr.font.color.rgb = _hex(ctx.body_color)

        return 'title'

    # ─────────────────────────────────────────────
    # CONTENT SLIDE
    # ─────────────────────────────────────────────

    def _content_slide(self, slide, slide_data, ctx: DesignContext):
        visual_type = slide_data.get('visual_type', 'none')
        has_visual  = visual_type not in ('none', None, '')

        # ── Text zone ──────────────────────────────────────────────
        # If visual: text takes left portion, visual takes right
        # If no visual: text spans full safe width
        if has_visual:
            text_w = ctx.w * 0.56
        else:
            text_w = ctx.w

        text_x = ctx.x0
        title_y = ctx.y0

        # ── Slide title ────────────────────────────────────────────
        title_h = Inches(0.85)
        ts = slide.shapes.add_textbox(text_x, title_y, text_w, title_h)
        tf = ts.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.alignment = ctx.title_align
        run = p.runs[0]
        run.font.name  = ctx.font_title
        run.font.size  = Pt(ctx.title_pt)
        run.font.bold  = ctx.category not in ('Minimal', 'Creative')
        run.font.color.rgb = _hex(ctx.title_color)

        # ── Divider under title ────────────────────────────────────
        if ctx.show_divider:
            div_y = title_y + title_h + Inches(0.05)
            dw = text_w * 0.35
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, text_x, div_y, dw, Inches(0.03))
            div.fill.solid()
            div.fill.fore_color.rgb = _hex(ctx.accent)
            div.line.fill.background()
            bullet_y = div_y + Inches(0.12)
        else:
            bullet_y = title_y + title_h + Inches(0.15)

        # ── Bullets ────────────────────────────────────────────────
        bullets = slide_data.get('bullets', [])
        n = len(bullets)

        # Adaptive sizing: more bullets → smaller font + tighter spacing
        if n <= 3:
            fsize, spacing, lspacing = Pt(ctx.body_pt_base + 1), Pt(12), ctx.line_spacing
        elif n <= 5:
            fsize, spacing, lspacing = Pt(ctx.body_pt_base),     Pt(9),  ctx.line_spacing
        else:
            fsize, spacing, lspacing = Pt(ctx.body_pt_base - 1), Pt(6),  1.25

        bullet_h = ctx.y1 - bullet_y
        cs = slide.shapes.add_textbox(text_x, bullet_y, text_w, bullet_h)
        ctf = cs.text_frame
        ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.TOP

        for i, bullet in enumerate(bullets[:6]):
            para = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            text = str(bullet).strip()
            if not text.startswith('•'):
                text = f'• {text}'
            para.text = text
            para.space_after  = spacing
            para.line_spacing = lspacing
            brun = para.runs[0]
            brun.font.name  = ctx.font_body
            brun.font.size  = fsize
            brun.font.bold  = False
            brun.font.color.rgb = _hex(ctx.body_color)

        return 'content'

    # ─────────────────────────────────────────────
    # VISUAL PLACEMENT
    # ─────────────────────────────────────────────

    def add_single_beautiful_visual(self, slide, slide_data, design_style,
                                    layout_info, visual_type, slide_index=0,
                                    total_slides=10):
        """Add ONE visual in the right-side safe zone. Returns True if added."""
        try:
            vtype = slide_data.get('visual_type', 'none')
            if vtype in ('none', None, ''):
                return False

            ctx = DesignContext(design_style)

            # Visual zone: right portion of safe area
            vis_x = ctx.x0 + ctx.w * 0.60
            vis_y = ctx.y0 + Inches(0.9)   # below title
            vis_w = ctx.x1 - vis_x
            vis_h = ctx.y1 - vis_y

            # Chart colour theme matches design
            theme = 'dark' if ctx.dark_bg else 'light'

            if vtype == 'image':
                return self._add_image(slide, slide_data, vis_x, vis_y, vis_w, vis_h)
            elif vtype == 'chart':
                return self._add_chart(slide, slide_data, vis_x, vis_y, vis_w, vis_h,
                                       'bar', theme, ctx.accent)
            elif vtype == 'pie':
                return self._add_chart(slide, slide_data, vis_x, vis_y, vis_w, vis_h,
                                       'pie', theme, ctx.accent)
            elif vtype == 'table':
                return self._add_table(slide, slide_data, vis_x, vis_y, vis_w, vis_h,
                                       theme, ctx.accent)
            return False

        except Exception as e:
            print(f"[VISUAL] Error: {e}")
            return False

    # ─────────────────────────────────────────────
    # VISUAL HELPERS
    # ─────────────────────────────────────────────

    def _add_image(self, slide, slide_data, x, y, w, h):
        try:
            from image_api_service import image_api
            img_data = image_api.get_image_for_slide(slide_data)
            if img_data and 'url' in img_data:
                img_bytes = image_api.download_image(img_data['url'])
                if img_bytes:
                    slide.shapes.add_picture(io.BytesIO(img_bytes), x, y, w, h)
                    print(f"[VISUAL] ✓ Image added")
                    return True
            return False
        except Exception as e:
            print(f"[VISUAL] Image error: {e}")
            return False

    def _add_chart(self, slide, slide_data, x, y, w, h,
                   chart_type, theme, accent_hex):
        try:
            from chart_service import chart_service
            data  = chart_service.generate_chart_data(slide_data, chart_type)
            img_b = chart_service.create_chart_image(data, theme, accent_hex)
            if img_b:
                slide.shapes.add_picture(io.BytesIO(img_b), x, y, w, h)
                print(f"[VISUAL] ✓ {chart_type} chart added ({theme} theme)")
                return True
            return False
        except Exception as e:
            print(f"[VISUAL] Chart error: {e}")
            return False

    def _add_table(self, slide, slide_data, x, y, w, h, theme, accent_hex):
        try:
            from chart_service import chart_service
            df    = chart_service.create_table_data(slide_data)
            img_b = chart_service.create_table_image(df, theme, accent_hex)
            if img_b:
                slide.shapes.add_picture(io.BytesIO(img_b), x, y, w, h)
                print(f"[VISUAL] ✓ Table added ({theme} theme)")
                return True
            return False
        except Exception as e:
            print(f"[VISUAL] Table error: {e}")
            return False

    # ─────────────────────────────────────────────
    # UTILITIES
    # ─────────────────────────────────────────────

    def _clear_slide(self, slide):
        for shape in list(slide.shapes):
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception:
                pass


# Global instance
beautiful_system = BeautifulSimpleSystem()
