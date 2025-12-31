from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# Beautiful theme definitions with attractive gradients and colors
PROFESSIONAL_THEMES = {
    "modern_blue": {
        "background": {
            "type": "gradient",
            "colors": ["#E0F2FE", "#BAE6FD", "#7DD3FC"],  # Beautiful blue gradient
            "direction": "vertical"  # top to bottom
        },
        "colors": {
            "title": "#0C4A6E",      # Deep blue for contrast
            "body": "#1E3A8A",       # Rich blue
            "accent": "#0284C7",     # Bright blue accent
            "background": "#E0F2FE"
        },
        "fonts": {
            "title": "Segoe UI",
            "subtitle": "Segoe UI Semibold",
            "body": "Segoe UI"
        },
        "spacing": {
            "title_top": 0.5,
            "content_top": 2.0,
            "margin": 0.8
        }
    },
    "academic_green": {
        "background": {
            "type": "gradient",
            "colors": ["#ECFDF5", "#D1FAE5", "#A7F3D0"],  # Soft green gradient
            "direction": "vertical"
        },
        "colors": {
            "title": "#065F46",      # Deep green
            "body": "#047857",       # Medium green
            "accent": "#10B981",     # Emerald accent
            "background": "#ECFDF5"
        },
        "fonts": {
            "title": "Segoe UI",
            "subtitle": "Segoe UI Semibold",
            "body": "Segoe UI"
        },
        "spacing": {
            "title_top": 0.5,
            "content_top": 2.0,
            "margin": 0.8
        }
    },
    "dark_mode": {
        "background": {
            "type": "gradient",
            "colors": ["#1F2937", "#374151", "#4B5563"],  # Dark gradient
            "direction": "diagonal"  # diagonal for depth
        },
        "colors": {
            "title": "#F9FAFB",      # Light text
            "body": "#E5E7EB",       # Light gray text
            "accent": "#60A5FA",     # Light blue accent
            "background": "#1F2937"
        },
        "fonts": {
            "title": "Segoe UI",
            "subtitle": "Segoe UI Semibold",
            "body": "Segoe UI"
        },
        "spacing": {
            "title_top": 0.5,
            "content_top": 2.0,
            "margin": 0.8
        }
    },
    "creative_orange": {
        "background": {
            "type": "gradient",
            "colors": ["#FFF7ED", "#FFEDD5", "#FED7AA"],  # Warm orange gradient
            "direction": "vertical"
        },
        "colors": {
            "title": "#9A3412",      # Deep orange
            "body": "#C2410C",       # Rich orange
            "accent": "#FB923C",     # Bright orange
            "background": "#FFF7ED"
        },
        "fonts": {
            "title": "Segoe UI",
            "subtitle": "Segoe UI Semibold",
            "body": "Segoe UI"
        },
        "spacing": {
            "title_top": 0.5,
            "content_top": 2.0,
            "margin": 0.8
        }
    },
    "minimal_white": {
        "background": {
            "type": "gradient",
            "colors": ["#FFFFFF", "#F9FAFB", "#F3F4F6"],  # Subtle white to gray
            "direction": "vertical"
        },
        "colors": {
            "title": "#111827",      # Near-black
            "body": "#374151",       # Medium gray
            "accent": "#6B7280",     # Subtle gray accent
            "background": "#FFFFFF"
        },
        "fonts": {
            "title": "Segoe UI",
            "subtitle": "Segoe UI Semibold",
            "body": "Segoe UI"
        },
        "spacing": {
            "title_top": 0.5,
            "content_top": 2.0,
            "margin": 0.8
        }
    }
}

def get_theme(theme_name: str) -> dict:
    """Get theme configuration by name"""
    return PROFESSIONAL_THEMES.get(theme_name, PROFESSIONAL_THEMES["minimal_white"])


def apply_theme(slide, theme):
    """
    Apply beautiful theme to slide with gradients and enhanced styling
    Supports both theme dict and theme name string
    """
    # Handle both theme dict and theme name string
    if isinstance(theme, str):
        theme = get_theme(theme)
    
    bg = slide.background
    fill = bg.fill

    # Beautiful background styling with gradients
    if theme["background"]["type"] == "gradient":
        fill.gradient()
        gradient_stops = fill.gradient_stops
        colors = theme["background"]["colors"]
        direction = theme["background"].get("direction", "vertical")
        
        # Modify existing gradient stops (gradient() creates 2 by default)
        # Set colors for the two default stops
        if len(colors) >= 2 and len(gradient_stops) >= 2:
            # First stop - use first color
            gradient_stops[0].color.rgb = hex_to_rgb(colors[0])
            gradient_stops[0].position = 0.0
            
            # Second stop - use last color
            gradient_stops[1].color.rgb = hex_to_rgb(colors[-1])
            gradient_stops[1].position = 1.0
        
        # Set gradient direction (90 = top-bottom, 0 = left-right)
        if direction == "diagonal":
            fill.gradient_angle = 45.0
        elif direction == "horizontal":
            fill.gradient_angle = 0.0
        else:  # vertical (default)
            fill.gradient_angle = 90.0
            
    elif theme["background"]["type"] == "solid":
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(theme["background"]["colors"][0])

    # Title styling with improved hierarchy
    if slide.shapes.title:
        title = slide.shapes.title
        title_frame = title.text_frame
        
        # Position and margins
        title.top = Inches(theme["spacing"]["title_top"])
        title_frame.margin_left = Inches(theme["spacing"]["margin"])
        title_frame.margin_right = Inches(theme["spacing"]["margin"])
        title_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Apply theme colors and fonts
        for p in title_frame.paragraphs:
            p.font.name = theme["fonts"]["title"]
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(theme["colors"]["title"])
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)

    # Body text styling
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # Set margins for better whitespace
        tf.margin_left = Inches(theme["spacing"]["margin"])
        tf.margin_right = Inches(theme["spacing"]["margin"])
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        tf.word_wrap = True
        
        # Apply body styling
        for p in tf.paragraphs:
            p.font.name = theme["fonts"]["body"]
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(theme["colors"]["body"])
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
            p.line_spacing = 1.2
